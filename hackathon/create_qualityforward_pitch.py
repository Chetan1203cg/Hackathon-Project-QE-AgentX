"""Populate the supplied QualityForward pitch-deck template for QE AgentX."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt


TEMPLATE_PATH = Path(
    r"C:\Chetan Raut\Projects\Hackathon 2026\Submission Data\QualityForward-Pitch-Deck-Template.pptx"
)
OUTPUT_PATH = Path(__file__).with_name("QE_AgentX_QualityForward_Pitch_Deck.pptx")


def replace_text(slide, starts_with, replacement, font_size=None, bold=None, align=None):
    """Replace a template prompt while retaining the template shape and layout."""
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text.strip().startswith(starts_with):
            continue

        frame = shape.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = replacement
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        if font_size is not None:
            paragraph.font.size = Pt(font_size)
        if bold is not None:
            paragraph.font.bold = bold
        if align is not None:
            paragraph.alignment = align
        return
    raise ValueError(f"Template text beginning with {starts_with!r} was not found.")


def set_text_frame(slide, starts_with, replacement, font_size=None, align=None):
    """Replace a multi-line prompt using one paragraph per supplied line."""
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text.strip().startswith(starts_with):
            continue

        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for index, line in enumerate(replacement.split("\n")):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(3)
            if font_size is not None:
                paragraph.font.size = Pt(font_size)
            if align is not None:
                paragraph.alignment = align
        return
    raise ValueError(f"Template text beginning with {starts_with!r} was not found.")


def remove_instruction_slide(presentation):
    """Remove the template's first, instructional slide."""
    slide_id = presentation.slides._sldIdLst[0]
    presentation.part.drop_rel(slide_id.rId)
    del presentation.slides._sldIdLst[0]


def populate_deck():
    presentation = Presentation(TEMPLATE_PATH)
    remove_instruction_slide(presentation)

    # Slide 1: title
    slide = presentation.slides[0]
    replace_text(slide, "[ TRACK:", "TRACK: Autonomous QE for Operations", 11, True)
    replace_text(slide, "[ Your Solution Name", "QE AgentX", 30, True)
    replace_text(
        slide,
        "[ One-line promise:",
        "AI-native Quality Engineering Control Tower for faster, traceable, explainable release decisions.",
        17,
    )
    replace_text(slide, "[ Team name", "VW QE AgentX | Chetan Raut | Idea", 11)

    # Slide 2: problem and user
    slide = presentation.slides[1]
    set_text_frame(
        slide,
        "Who is the user?",
        "Release Manager and QA Lead\nTurn Jira stories, test evidence, and defects into a trusted release decision.\nToday: manual RTM, coverage analysis, and release reporting every sprint.",
        15,
    )
    set_text_frame(
        slide,
        "Cost of the status quo",
        "8-12 hours of QA reporting, traceability, and release governance per release cycle.\nDisconnected Jira, Excel, Confluence, emails, and evidence repositories create gaps and duplicate work.\nPOV: QA leads need one reliable quality view because scattered evidence delays confident Go/No-Go decisions.",
        15,
    )

    # Slide 3: solution and agent flow
    slide = presentation.slides[2]
    set_text_frame(
        slide,
        "What is the solution,",
        "QE AgentX is a Test Asset Intelligence Platform that connects requirements, testing, execution, defects, and release governance.\nAI interprets requirements and acceptance criteria, generates test assets, maintains traceability, detects quality risks, and explains a release recommendation.\nSpecialized agents: Requirement Intelligence, Test Design, Traceability, Knowledge, Test Debt, Automation, Sprint Analytics, and Release Readiness.\nHuman review approves critical test assets and release decisions.",
        13,
    )
    set_text_frame(
        slide,
        "PASTE YOUR ARCHITECTURE",
        "Jira Story + Acceptance Criteria + Execution Evidence\n↓\nRequirement Intelligence → Test Design → RTM & Knowledge Store\n↓\nExecution & Defect Intelligence → Sprint Analytics\n↓\nExplainable GO / GO WITH RISKS / NO-GO",
        15,
        PP_ALIGN.CENTER,
    )
    replace_text(
        slide,
        "A simple boxes-and-arrows",
        "Inputs become traceable test assets, quality insights, and a governed release recommendation.",
        12,
    )

    # Slide 4: demo flow
    slide = presentation.slides[3]
    set_text_frame(
        slide,
        "SCREENSHOT / MOCK",
        "LIVE MVP DEMO\n\n1. Select a Jira story and its acceptance criteria\n2. Review generated scenarios, test cases, and RTM links\n3. Load execution results and defect evidence\n4. View coverage, risks, and the release-readiness recommendation",
        17,
        PP_ALIGN.CENTER,
    )
    set_text_frame(
        slide,
        "Step 1 — user does…",
        "Step 1 - QA lead selects a Jira requirement.\nStep 2 - QE AgentX builds tests, RTM links, and coverage evidence.\nStep 3 - Execution and defect signals produce an explainable release view.\nBuilt for MVP: story ingestion, test generation, RTM, dashboard, and recommendation.",
        14,
    )

    # Slide 5: measurable impact
    slide = presentation.slides[4]
    replace_text(
        slide,
        "Also note:",
        "Measured per release cycle: baseline manual QA governance effort versus the time to generate traceability, quality insights, and a recommendation. Target is validated in a controlled MVP pilot.",
        13,
    )
    metric_labels = [
        ("8-12 hrs", "Manual reporting, RTM, and release governance today"),
        ("1-2 hrs", "Target effort with AI-driven automation"),
        (">90%", "Target reduction in manual governance effort"),
    ]
    metric_shapes = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.width > 1000000
    ]
    for shape, (metric, label) in zip(metric_shapes[:3], metric_labels):
        frame = shape.text_frame
        frame.clear()
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        top = frame.paragraphs[0]
        top.text = metric
        top.font.size = Pt(26)
        top.font.bold = True
        top.alignment = PP_ALIGN.CENTER
        bottom = frame.add_paragraph()
        bottom.text = label
        bottom.font.size = Pt(11)
        bottom.alignment = PP_ALIGN.CENTER

    # Slide 6: differentiation and responsible AI
    slide = presentation.slides[5]
    set_text_frame(
        slide,
        "What exists today",
        "Point tools generate tests or store evidence; manual RTMs and release assessments remain disconnected.\nQE AgentX links Story → Acceptance Criteria → Test Case → Execution → Defect in one quality-intelligence loop.\nAI is essential because it synthesizes dispersed evidence, finds coverage and reuse gaps, and explains the decision.\nReusable across Jira-based Agile teams without requiring a full test-management-system rollout.",
        14,
    )
    set_text_frame(
        slide,
        "Key AI risks",
        "Risk: inaccurate generated assets, unsupported conclusions, and sensitive delivery data.\nGuardrails: human approval for critical outputs; source-linked evidence; deterministic rule thresholds; role-based access; local demo data.\nValidation: review generated cases and RTM links against acceptance criteria; reconcile decision inputs with execution and defect records.\nRelease recommendation is decision support, never an autonomous production sign-off.",
        14,
    )

    # Slide 7: feasibility and roadmap
    slide = presentation.slides[6]
    set_text_frame(
        slide,
        "Smallest slice",
        "MVP: Jira story → test scenarios/test cases → RTM → coverage and execution view → explainable release recommendation.\nDemo uses sample Jira stories, acceptance criteria, manual/automation results, and defects.\nDependencies: Jira API access and representative QE data.",
        13,
    )
    set_text_frame(
        slide,
        "Languages / frameworks",
        "Python, FastAPI, SQLite, Streamlit, Jira APIs, VS Code, GitHub Copilot.\nAI: GitHub Copilot Chat and Microsoft Copilot; Azure OpenAI or enterprise AI service for MVP validation.\nData: Jira stories, acceptance criteria, sprint/release metadata, execution results, and defects.",
        13,
    )
    set_text_frame(
        slide,
        "Effort / team / timeline",
        "Hackathon MVP: Chetan Raut / VW QE AgentX.\nPilot: validate two sprints with KPI tracking for cycle time, traceability completeness, coverage confidence, and defect leakage.\nV2: connectors, reusable test-asset intelligence, predictive risk signals, and broader team rollout.",
        13,
    )

    # Slide 8: ask and team
    slide = presentation.slides[7]
    set_text_frame(
        slide,
        "What support, data, access",
        "Jira API and sample QE-data access; Azure OpenAI or enterprise AI experimentation access; and mentoring on agentic architecture, RTM governance, KPI design, and demo storytelling.\nDecision requested: support a controlled two-sprint MVP pilot.\nNext step: validate the measurable effort reduction and release-confidence outcomes with real delivery data.",
        15,
    )
    set_text_frame(
        slide,
        "Team name & lead",
        "VW QE AgentX\nLead: Chetan Raut\nCapgemini\n\nAI-native quality intelligence for confident releases.",
        17,
    )
    replace_text(
        slide,
        "Close on the promise",
        "QE AgentX turns scattered quality evidence into traceable test intelligence and explainable release decisions.",
        14,
        True,
        PP_ALIGN.CENTER,
    )

    presentation.save(OUTPUT_PATH)
    print(f"Created {OUTPUT_PATH}")


if __name__ == "__main__":
    populate_deck()